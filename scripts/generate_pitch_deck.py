from __future__ import annotations

import json
import math
import shutil
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from textwrap import wrap

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.colors import HexColor
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas


ROOT = Path("/root/projetos/quarry")
OUT_DIR = ROOT / "docs" / "pitch"
SCREENSHOT = ROOT / "artifacts" / "public-deploy" / "quarry-public-demo-current.png"
SECRL_JSON = ROOT / "docs" / "benchmarks" / "secrl-results.json"
SECRL_MD = ROOT / "docs" / "benchmarks" / "microsoft-secrl-results.md"


COLORS = {
    "bg": "0B0F14",
    "panel": "151A21",
    "panel2": "1E252E",
    "red": "8B0000",
    "red2": "B32635",
    "rust": "A8482E",
    "white": "F7F4EF",
    "muted": "A8B0BA",
    "line": "343B45",
    "gold": "D9A441",
}


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def hx(hex_value: str) -> HexColor:
    return HexColor("#" + hex_value.strip("#"))


@dataclass
class SlideSpec:
    title: str
    kicker: str = ""
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    metrics: list[tuple[str, str, str]] = field(default_factory=list)
    source: str = ""
    kind: str = "standard"


def load_secrl_summary() -> dict:
    data = json.loads(SECRL_JSON.read_text())
    results = data.get("results", [])
    correct = sum(1 for r in results if r.get("correct"))
    total_tokens = sum(int(r.get("input_tokens") or 0) + int(r.get("output_tokens") or 0) for r in results)
    total_cost = sum(float(r.get("cost_usd") or 0.0) for r in results)
    avg_time = sum(float(r.get("time_seconds") or 0.0) for r in results) / max(len(results), 1)
    return {
        "generated_at": data.get("generated_at", ""),
        "model": data.get("model", "gpt-4o"),
        "eval_model": data.get("eval_model", "gpt-4o"),
        "temperature": data.get("temperature", 0),
        "max_steps": data.get("max_steps", 25),
        "limit_per_incident": data.get("limit_per_incident", 1),
        "scenarios": len(results),
        "correct": correct,
        "accuracy": correct / max(len(results), 1),
        "total_tokens": total_tokens,
        "total_cost": total_cost,
        "avg_time": avg_time,
    }


SECRL = load_secrl_summary()


SLIDES = [
    SlideSpec(
        kind="cover",
        kicker="Confidential discussion deck | May 2026",
        title="Quarry — Autonomous Threat Hunting for the Post-Mythos Era",
        subtitle="A governed agentic investigation platform for financial-sector cyber defense.",
        source="Demo: https://quarry.12brain.org",
    ),
    SlideSpec(
        kind="date",
        kicker="The Mythos Moment",
        title="April 7, 2026 changed the threat-hunting clock.",
        subtitle="Frontier agents are now credible enough to compress vulnerability discovery, exploitation, and defensive response cycles.",
        bullets=[
            "Anthropic disclosed Claude Mythos Preview and Project Glasswing, explicitly framing the launch as a security watershed.",
            "The Federal Reserve later cited Mythos at an FSOC AI / cybersecurity roundtable as evidence of rapidly evolving capability.",
            "The implication for financial institutions is operational: investigation workflows must move from analyst-hours to machine-minutes while remaining auditable.",
        ],
        source="Sources: Anthropic Red Team, Apr 7 2026; Federal Reserve Bowman speech, May 1 2026.",
    ),
    SlideSpec(
        kind="window",
        kicker="The Window",
        title="The next 6-12 months are the adoption window.",
        subtitle="Regulators are turning AI-cyber guidance into operating expectations while attackers gain agentic leverage first.",
        bullets=[
            "Treasury launched public-private AI cybersecurity resources for financial services in 2026.",
            "Treasury's earlier AI cybersecurity report identified immediate operational risk, cybersecurity, and fraud challenges.",
            "The buyer will not want a chatbot. The buyer will want governed autonomy: provenance, cost control, reproducibility, and policy guardrails.",
        ],
        source="Sources: U.S. Treasury AI cybersecurity initiative, Feb 18 2026; Treasury AI cybersecurity report, Mar 27 2024.",
    ),
    SlideSpec(
        kind="gap",
        kicker="The Gap",
        title="Hunters cannot scale linearly with telemetry.",
        subtitle="The industry has more alerts, more tools, more data, and fewer people who can connect it all under time pressure.",
        bullets=[
            "SIEM / EDR / cloud logs are searchable, but the investigation work is still mostly human orchestration.",
            "Senior hunters are the bottleneck: hypothesis generation, query selection, evidence correlation, and report writing.",
            "Static detections miss the operating need: turning a vague brief into a defensible investigation trail.",
        ],
        source="Sources: Microsoft SecRL benchmark framing; Quarry architecture docs.",
    ),
    SlideSpec(
        kind="solution",
        kicker="The Solution",
        title="Quarry turns a brief into governed investigation work.",
        subtitle="An agentic hunting platform that decomposes the brief, queries telemetry, correlates evidence, and writes the report with a ledger.",
        source="Source: Quarry product architecture and demo build.",
    ),
    SlideSpec(
        kind="screenshot",
        kicker="How It Works",
        title="The analyst sees a live investigation, not a black box.",
        subtitle="Brief, hypotheses, data queries, findings, graph, MITRE coverage, cost, and report generation are visible in one flow.",
        bullets=[
            "Input: natural-language investigation brief.",
            "Process: hypotheses -> connector queries -> evidence ledger -> report.",
            "Output: analyst-ready narrative with supporting artifacts.",
        ],
        source="Source: Quarry public demo screenshot captured May 16 2026.",
    ),
    SlideSpec(
        kind="proof",
        kicker="The Proof",
        title="SecRL smoke run: 8/8 scenarios correct.",
        subtitle="A pitch-grade regression run against Microsoft's cyber threat investigation benchmark, using the Quarry SQL investigation adapter.",
        metrics=[
            ("8/8", "correct scenarios", "one question per incident"),
            (f"{SECRL['avg_time']:.2f}s", "avg time", "brief to answer"),
            (f"{SECRL['total_tokens']:,}", "tokens", "total run"),
            (f"${SECRL['total_cost']:.3f}", "cost", "estimated total"),
        ],
        bullets=[
            f"Model: {SECRL['model']}; evaluator: {SECRL['eval_model']}; temperature: {SECRL['temperature']}; max steps: {SECRL['max_steps']}.",
            "Scope note: this is an 8-scenario pitch smoke run, not the full official question-set claim.",
        ],
        source="Sources: docs/benchmarks/microsoft-secrl-results.md; microsoft/SecRL GitHub.",
    ),
    SlideSpec(
        kind="beachhead",
        kicker="The Beachhead",
        title="Brazilian fintechs are the first wedge.",
        subtitle="Pix created a national-scale real-time payment environment where fraud, identity, mobile, and infrastructure signals converge.",
        metrics=[
            ("+170M", "individual Pix users", "BCB public stats"),
            ("+7B", "Pix transactions", "January 2026"),
            ("+R$3T", "monthly volume", "October 2025"),
        ],
        bullets=[
            "Demo persona: FinPlay Pagamentos, a fictitious Sao Paulo fintech with 50,000 active customers.",
            "Synthetic 30-day dataset covers Pix, mobile banking auth, Open Finance, boleto creation, and injected fraud campaigns.",
            "Initial use case: Pix fraud and account-takeover investigations that require fast correlation and an auditable report.",
        ],
        source="Sources: Banco Central do Brasil Pix statistics; CARD-009 synthetic dataset.",
    ),
    SlideSpec(
        kind="team",
        kicker="The Team",
        title="Founder-led product, partner-scaled distribution.",
        subtitle="The right early team pairs Brazilian market access with global enterprise-security credibility.",
        bullets=[
            "Jose / Increase Trainer: product vision, Brazilian market access, and hands-on operator discovery.",
            "Meta co-founder partner: enterprise credibility, security-network leverage, and strategic distribution.",
            "Operating model: narrow wedge, visible demo, measurable benchmark, then regulated pilot.",
        ],
        source="Source: founder / partner operating thesis.",
    ),
    SlideSpec(
        kind="roadmap",
        kicker="The Roadmap",
        title="BR -> US -> Global, with evidence gates at each step.",
        subtitle="The roadmap protects credibility by pairing expansion with benchmark, security, and deployment milestones.",
        source="Source: Quarry product roadmap, May 2026.",
    ),
    SlideSpec(
        kind="ask",
        kicker="The Ask",
        title="Structure the partnership now.",
        subtitle="The near-term ask is not broad fundraising. It is a focused partnership discussion tied to diligence and pilot execution.",
        bullets=[
            "Align on co-founder / advisor / equity structure and operating responsibilities.",
            "Run 30-day technical diligence: public demo, architecture, SecRL harness, and synthetic fintech dataset.",
            "Define 60-day pilot path with one Brazilian fintech or MSSP and one enterprise-security design partner.",
            "Decision point: strategic partnership, seed vehicle, or customer-led buildout.",
        ],
        source="Source: proposed partnership timeline.",
    ),
    SlideSpec(
        kind="close",
        kicker="Closing",
        title="Human-speed teams are entering machine-speed investigations.",
        subtitle="Quarry is the control plane that keeps autonomous hunting useful, governed, and defensible.",
        bullets=[
            "Live demo: https://quarry.12brain.org",
            "Backup artifacts: docs/pitch/quarry-v1.pdf and docs/pitch/quarry-v1.pptx",
            "Next steps: partner diligence, full SecRL question-set run, first regulated pilot.",
        ],
        source="Sources: Quarry demo and source notes in docs/pitch/quarry-v1-sources.md.",
    ),
]


def set_fill(shape, color: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 20,
             color: str = "white", bold: bool = False, font: str = "Inter",
             align=PP_ALIGN.LEFT, line_spacing: float | None = None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color])
    return box


def add_multiline(slide, lines: list[str], x: float, y: float, w: float, h: float,
                  size: int = 18, color: str = "white", bullet: bool = False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        run = p.add_run()
        run.text = ("• " if bullet else "") + line
        run.font.name = "Inter"
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(COLORS[color])
    return box


def add_shell(slide, spec: SlideSpec):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(COLORS["bg"])
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.08)).fill.solid()
    top = slide.shapes[-1]
    top.fill.fore_color.rgb = rgb(COLORS["red"])
    top.line.color.rgb = rgb(COLORS["red"])
    add_logo(slide, 0.55, 0.34, 0.38)
    if spec.kicker:
        add_text(slide, spec.kicker.upper(), 1.05, 0.38, 6.0, 0.22, 8, "muted", True)
    if spec.source:
        add_text(slide, spec.source, 0.55, 7.06, 12.3, 0.26, 7, "muted")


def add_logo(slide, x: float, y: float, size: float):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.HEXAGON,
        Inches(x),
        Inches(y),
        Inches(size),
        Inches(size),
    )
    set_fill(shape, COLORS["red"])
    add_text(slide, "Q", x + 0.105 * size, y + 0.045 * size, size * 0.8, size * 0.8, 13, "white", True, align=PP_ALIGN.CENTER)


def add_metric_card(slide, x: float, y: float, w: float, h: float, value: str, label: str, sub: str):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(card, COLORS["panel"])
    card.line.color.rgb = rgb(COLORS["line"])
    add_text(slide, value, x + 0.18, y + 0.18, w - 0.36, 0.42, 26, "white", True)
    add_text(slide, label.upper(), x + 0.2, y + 0.78, w - 0.4, 0.26, 8, "gold", True)
    add_text(slide, sub, x + 0.2, y + 1.12, w - 0.4, 0.34, 8, "muted")


def add_pipeline(slide, x: float, y: float):
    labels = ["Brief", "Hypotheses", "Federated\nSearch", "Evidence\nLedger", "Report"]
    for i, label in enumerate(labels):
        bx = x + i * 2.35
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(bx), Inches(y), Inches(1.7), Inches(0.82))
        set_fill(box, COLORS["panel2"] if i != 3 else COLORS["red"])
        box.line.color.rgb = rgb(COLORS["line"])
        add_text(slide, label, bx + 0.12, y + 0.16, 1.46, 0.5, 14, "white", True, align=PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            add_text(slide, "->", bx + 1.82, y + 0.25, 0.45, 0.3, 14, "muted", True, align=PP_ALIGN.CENTER)


def create_pptx(path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, spec in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        add_shell(slide, spec)
        if spec.kind == "cover":
            add_text(slide, "QUARRY", 0.7, 1.65, 7.8, 0.9, 58, "white", True, font="Aptos Display")
            add_text(slide, "Autonomous Threat Hunting for the Post-Mythos Era", 0.78, 2.62, 8.4, 0.5, 23, "muted")
            add_text(slide, spec.subtitle, 0.8, 3.45, 6.3, 0.9, 19, "white")
            add_text(slide, "Confidential", 0.83, 5.88, 2.2, 0.3, 10, "gold", True)
            add_text(slide, "Prepared for strategic partner discussion", 0.83, 6.23, 5.2, 0.3, 11, "muted")
            for i, radius in enumerate([4.0, 3.0, 2.0]):
                s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, Inches(8.15 + i * 0.42), Inches(1.45 + i * 0.2), Inches(radius), Inches(radius))
                s.fill.solid()
                s.fill.fore_color.rgb = rgb(COLORS["panel2"] if i % 2 else COLORS["red"])
                s.fill.transparency = 15 + i * 20
                s.line.color.rgb = rgb(COLORS["line"])
            add_text(slide, "Q", 9.38, 2.58, 1.4, 1.2, 64, "white", True, align=PP_ALIGN.CENTER)
        elif spec.kind == "date":
            add_text(slide, "APR", 0.83, 1.2, 1.8, 0.35, 17, "gold", True)
            add_text(slide, "07", 0.72, 1.48, 2.5, 1.2, 82, "red2", True)
            add_text(slide, "2026", 0.87, 2.6, 2.0, 0.35, 19, "muted", True)
            add_text(slide, spec.title, 3.2, 1.2, 8.8, 0.95, 35, "white", True)
            add_text(slide, spec.subtitle, 3.25, 2.24, 8.6, 0.7, 17, "muted")
            add_multiline(slide, spec.bullets, 3.3, 3.2, 8.7, 2.3, 17, "white", True)
        elif spec.kind == "window":
            add_text(slide, spec.title, 0.8, 1.05, 10.7, 0.7, 34, "white", True)
            add_text(slide, "6-12", 0.82, 2.0, 3.2, 0.95, 64, "red2", True)
            add_text(slide, "months", 3.1, 2.34, 1.8, 0.4, 22, "muted", True)
            add_text(slide, spec.subtitle, 5.0, 2.05, 6.5, 0.75, 18, "white")
            add_multiline(slide, spec.bullets, 0.95, 3.3, 10.9, 2.2, 18, "white", True)
            add_text(slide, "Positioning rule: auditable autonomy beats unsupervised autonomy.", 0.95, 5.85, 8.5, 0.38, 15, "gold", True)
        elif spec.kind == "gap":
            add_text(slide, spec.title, 0.8, 1.0, 10.8, 0.65, 34, "white", True)
            add_text(slide, spec.subtitle, 0.83, 1.85, 10.4, 0.5, 16, "muted")
            labels = [("Telemetry", "more volume"), ("Tools", "more surface"), ("Hunters", "same bottleneck"), ("Reports", "late evidence")]
            for i, (a, b) in enumerate(labels):
                add_metric_card(slide, 0.85 + i * 3.05, 2.75, 2.55, 1.55, a, b, "current SOC reality")
            add_multiline(slide, spec.bullets, 1.0, 4.85, 10.9, 1.35, 16, "white", True)
        elif spec.kind == "solution":
            add_text(slide, spec.title, 0.8, 1.0, 11.2, 0.7, 33, "white", True)
            add_text(slide, spec.subtitle, 0.82, 1.84, 10.9, 0.65, 17, "muted")
            add_pipeline(slide, 0.92, 3.05)
            callouts = [
                ("Governed", "Policy guardrails and read-only evaluation paths."),
                ("Traceable", "Every query and conclusion is written to the Investigation Ledger."),
                ("Measurable", "Time, tokens, cost, and accuracy become product metrics."),
            ]
            for i, (head, body) in enumerate(callouts):
                add_metric_card(slide, 1.05 + i * 4.0, 4.85, 3.45, 1.35, head, body, "design principle")
        elif spec.kind == "screenshot":
            add_text(slide, spec.title, 0.8, 0.95, 7.8, 0.55, 28, "white", True)
            add_text(slide, spec.subtitle, 0.82, 1.55, 7.4, 0.55, 13, "muted")
            if SCREENSHOT.exists():
                slide.shapes.add_picture(str(SCREENSHOT), Inches(0.82), Inches(2.15), width=Inches(7.85))
            add_multiline(slide, spec.bullets, 9.05, 2.1, 3.35, 1.85, 15, "white", True)
            add_pipeline(slide, 1.0, 6.1)
        elif spec.kind == "proof":
            add_text(slide, spec.title, 0.8, 0.95, 10.0, 0.65, 34, "white", True)
            add_text(slide, spec.subtitle, 0.82, 1.76, 10.7, 0.55, 16, "muted")
            for i, m in enumerate(spec.metrics):
                add_metric_card(slide, 0.82 + i * 3.1, 2.75, 2.65, 1.55, *m)
            add_multiline(slide, spec.bullets, 0.95, 4.85, 10.8, 1.0, 16, "white", True)
            add_text(slide, "Conservative claim: early evidence of fit, not a completed official benchmark certification.", 0.95, 6.05, 10.8, 0.35, 13, "gold", True)
        elif spec.kind == "beachhead":
            add_text(slide, spec.title, 0.8, 0.95, 10.2, 0.65, 33, "white", True)
            add_text(slide, spec.subtitle, 0.82, 1.72, 10.9, 0.55, 16, "muted")
            for i, m in enumerate(spec.metrics):
                add_metric_card(slide, 0.82 + i * 4.08, 2.55, 3.45, 1.55, *m)
            add_multiline(slide, spec.bullets, 0.95, 4.55, 10.9, 1.55, 16, "white", True)
        elif spec.kind == "team":
            add_text(slide, spec.title, 0.8, 0.95, 10.3, 0.65, 33, "white", True)
            add_text(slide, spec.subtitle, 0.82, 1.72, 10.4, 0.5, 16, "muted")
            cards = [
                ("Jose", "Product vision\nBrazilian market access\nOperator discovery"),
                ("Meta partner", "Enterprise credibility\nSecurity network leverage\nStrategic distribution"),
                ("Operating model", "Demo first\nBenchmark visible\nPilot disciplined"),
            ]
            for i, (head, body) in enumerate(cards):
                add_metric_card(slide, 0.9 + i * 4.05, 2.7, 3.45, 2.2, head, body.replace("\n", " | "), "role")
            add_multiline(slide, spec.bullets, 1.05, 5.35, 10.8, 1.0, 15, "white", True)
        elif spec.kind == "roadmap":
            add_text(slide, spec.title, 0.8, 0.95, 10.5, 0.65, 33, "white", True)
            add_text(slide, spec.subtitle, 0.82, 1.72, 10.8, 0.5, 16, "muted")
            stages = [
                ("Q2 2026", "Public demo\nBR synthetic fintech\nSecRL smoke run"),
                ("Q3 2026", "Full SecRL run\nSOC connectors\nFirst BR pilot"),
                ("Q4 2026", "Compliance hardening\nUS finance beachhead\nPartner packaging"),
                ("H1 2027", "Global pilots\nChannel motion\nEnterprise controls"),
            ]
            y = 3.0
            for i, (head, body) in enumerate(stages):
                x = 0.95 + i * 3.05
                add_metric_card(slide, x, y, 2.65, 2.15, head, body.replace("\n", " | "), "milestone gate")
            add_text(slide, "Rule: every expansion step ships with a repeatable demo, benchmark evidence, and operational runbook.", 1.0, 5.85, 10.8, 0.38, 14, "gold", True)
        elif spec.kind == "ask":
            add_text(slide, spec.title, 0.8, 0.95, 10.5, 0.65, 34, "white", True)
            add_text(slide, spec.subtitle, 0.82, 1.72, 10.5, 0.5, 16, "muted")
            add_multiline(slide, spec.bullets, 0.95, 2.65, 10.8, 2.35, 18, "white", True)
            steps = [("Week 1", "alignment", "partnership terms"), ("Day 30", "diligence", "technical evidence"), ("Day 60", "pilot decision", "customer path")]
            for i, m in enumerate(steps):
                add_metric_card(slide, 1.0 + i * 4.05, 5.35, 3.25, 1.25, *m)
        elif spec.kind == "close":
            add_text(slide, spec.title, 0.82, 1.15, 10.7, 0.9, 40, "white", True)
            add_text(slide, spec.subtitle, 0.86, 2.28, 8.8, 0.65, 19, "muted")
            add_multiline(slide, spec.bullets, 0.95, 3.45, 9.8, 1.4, 18, "white", True)
            add_text(slide, "quarry.12brain.org", 0.95, 5.72, 5.4, 0.55, 27, "red2", True)
            add_logo(slide, 10.2, 4.55, 1.25)

        add_text(slide, f"{idx:02d}/12", 12.35, 7.06, 0.45, 0.22, 7, "muted", True, align=PP_ALIGN.RIGHT)

    prs.save(path)


def pdf_wrap(text: str, max_chars: int) -> list[str]:
    lines: list[str] = []
    for part in text.split("\n"):
        lines.extend(wrap(part, max_chars) if part else [""])
    return lines


def pdf_text(c, text: str, x: float, y: float, size: int = 18, color: str = "white",
             font: str = "Helvetica", leading: float | None = None, max_chars: int | None = None):
    c.setFillColor(hx(COLORS[color]))
    c.setFont(font, size)
    lines = pdf_wrap(text, max_chars) if max_chars else text.split("\n")
    if leading is None:
        leading = size * 1.25
    for i, line in enumerate(lines):
        c.drawString(x, y - i * leading, line)
    return y - len(lines) * leading


def pdf_bullets(c, bullets: list[str], x: float, y: float, size: int = 15, max_chars: int = 82):
    c.setFillColor(hx(COLORS["white"]))
    c.setFont("Helvetica", size)
    leading = size * 1.45
    for bullet in bullets:
        lines = pdf_wrap(bullet, max_chars)
        c.drawString(x, y, "• " + lines[0])
        y -= leading
        for line in lines[1:]:
            c.drawString(x + 14, y, line)
            y -= leading
        y -= 3
    return y


def pdf_card(c, x: float, y: float, w: float, h: float, value: str, label: str, sub: str):
    c.setFillColor(hx(COLORS["panel"]))
    c.setStrokeColor(hx(COLORS["line"]))
    c.roundRect(x, y, w, h, 8, fill=1, stroke=1)
    pdf_text(c, value, x + 14, y + h - 34, 24, "white", "Helvetica-Bold")
    pdf_text(c, label.upper(), x + 14, y + h - 58, 8, "gold", "Helvetica-Bold")
    pdf_text(c, sub, x + 14, y + 18, 8, "muted", max_chars=25)


def pdf_shell(c, spec: SlideSpec, idx: int):
    width, height = 960, 540
    c.setFillColor(hx(COLORS["bg"]))
    c.rect(0, 0, width, height, fill=1, stroke=0)
    c.setFillColor(hx(COLORS["red"]))
    c.rect(0, height - 6, width, 6, fill=1, stroke=0)
    c.setFillColor(hx(COLORS["red"]))
    c.roundRect(40, height - 52, 26, 26, 8, fill=1, stroke=0)
    pdf_text(c, "Q", 48, height - 46, 13, "white", "Helvetica-Bold")
    if spec.kicker:
        pdf_text(c, spec.kicker.upper(), 78, height - 43, 8, "muted", "Helvetica-Bold")
    if spec.source:
        pdf_text(c, spec.source, 40, 25, 7, "muted", max_chars=150)
    pdf_text(c, f"{idx:02d}/12", 890, 25, 7, "muted", "Helvetica-Bold")


def create_pdf(path: Path):
    c = canvas.Canvas(str(path), pagesize=(960, 540))
    for idx, spec in enumerate(SLIDES, start=1):
        pdf_shell(c, spec, idx)
        if spec.kind == "cover":
            pdf_text(c, "QUARRY", 58, 380, 56, "white", "Helvetica-Bold")
            pdf_text(c, "Autonomous Threat Hunting for the Post-Mythos Era", 62, 326, 22, "muted")
            pdf_text(c, spec.subtitle, 62, 260, 18, "white", max_chars=58)
            c.setFillColor(hx(COLORS["red"]))
            c.roundRect(620, 160, 230, 230, 38, fill=1, stroke=0)
            pdf_text(c, "Q", 700, 280, 72, "white", "Helvetica-Bold")
            pdf_text(c, "Confidential", 62, 92, 10, "gold", "Helvetica-Bold")
            pdf_text(c, "Prepared for strategic partner discussion", 62, 68, 10, "muted")
        elif spec.kind == "date":
            pdf_text(c, "APR", 62, 430, 17, "gold", "Helvetica-Bold")
            pdf_text(c, "07", 56, 340, 80, "red2", "Helvetica-Bold")
            pdf_text(c, "2026", 66, 302, 19, "muted", "Helvetica-Bold")
            pdf_text(c, spec.title, 230, 410, 32, "white", "Helvetica-Bold", max_chars=42)
            pdf_text(c, spec.subtitle, 232, 334, 15, "muted", max_chars=76)
            pdf_bullets(c, spec.bullets, 238, 270, 15, 72)
        elif spec.kind == "window":
            pdf_text(c, spec.title, 58, 414, 31, "white", "Helvetica-Bold")
            pdf_text(c, "6-12", 60, 330, 64, "red2", "Helvetica-Bold")
            pdf_text(c, "months", 225, 318, 22, "muted", "Helvetica-Bold")
            pdf_text(c, spec.subtitle, 360, 340, 17, "white", max_chars=55)
            pdf_bullets(c, spec.bullets, 72, 250, 15, 90)
            pdf_text(c, "Positioning rule: auditable autonomy beats unsupervised autonomy.", 72, 82, 13, "gold", "Helvetica-Bold")
        elif spec.kind == "gap":
            pdf_text(c, spec.title, 58, 416, 31, "white", "Helvetica-Bold")
            pdf_text(c, spec.subtitle, 60, 360, 15, "muted", max_chars=94)
            for i, (a, b) in enumerate([("Telemetry", "more volume"), ("Tools", "more surface"), ("Hunters", "same bottleneck"), ("Reports", "late evidence")]):
                pdf_card(c, 60 + i * 220, 215, 185, 96, a, b, "current SOC reality")
            pdf_bullets(c, spec.bullets, 72, 170, 14, 95)
        elif spec.kind == "solution":
            pdf_text(c, spec.title, 58, 410, 31, "white", "Helvetica-Bold")
            pdf_text(c, spec.subtitle, 60, 350, 15, "muted", max_chars=94)
            labels = ["Brief", "Hypotheses", "Federated Search", "Evidence Ledger", "Report"]
            for i, label in enumerate(labels):
                x = 70 + i * 170
                c.setFillColor(hx(COLORS["red"] if i == 3 else COLORS["panel2"]))
                c.roundRect(x, 250, 130, 50, 8, fill=1, stroke=0)
                pdf_text(c, label, x + 15, 270, 12, "white", "Helvetica-Bold")
                if i < 4:
                    pdf_text(c, "->", x + 142, 268, 13, "muted", "Helvetica-Bold")
            for i, (head, body) in enumerate([("Governed", "Policy guardrails and read-only evaluation paths."), ("Traceable", "Every query and conclusion goes to the ledger."), ("Measurable", "Time, tokens, cost, and accuracy are product metrics.")]):
                pdf_card(c, 75 + i * 280, 90, 240, 88, head, body, "design principle")
        elif spec.kind == "screenshot":
            pdf_text(c, spec.title, 58, 430, 26, "white", "Helvetica-Bold")
            pdf_text(c, spec.subtitle, 60, 390, 13, "muted", max_chars=115)
            if SCREENSHOT.exists():
                img = ImageReader(str(SCREENSHOT))
                iw, ih = Image.open(SCREENSHOT).size
                target_w = 430
                target_h = target_w * ih / iw
                c.drawImage(img, 60, 78, width=target_w, height=target_h, preserveAspectRatio=True, mask="auto")
            pdf_bullets(c, spec.bullets, 675, 320, 14, 36)
        elif spec.kind == "proof":
            pdf_text(c, spec.title, 58, 420, 31, "white", "Helvetica-Bold")
            pdf_text(c, spec.subtitle, 60, 362, 15, "muted", max_chars=92)
            for i, m in enumerate(spec.metrics):
                pdf_card(c, 60 + i * 220, 220, 185, 96, *m)
            pdf_bullets(c, spec.bullets, 72, 170, 14, 95)
            pdf_text(c, "Conservative claim: early evidence of fit, not a completed official benchmark certification.", 72, 72, 12, "gold", "Helvetica-Bold")
        elif spec.kind == "beachhead":
            pdf_text(c, spec.title, 58, 420, 31, "white", "Helvetica-Bold")
            pdf_text(c, spec.subtitle, 60, 362, 15, "muted", max_chars=94)
            for i, m in enumerate(spec.metrics):
                pdf_card(c, 60 + i * 290, 225, 245, 96, *m)
            pdf_bullets(c, spec.bullets, 72, 170, 14, 95)
        elif spec.kind == "team":
            pdf_text(c, spec.title, 58, 420, 31, "white", "Helvetica-Bold")
            pdf_text(c, spec.subtitle, 60, 362, 15, "muted", max_chars=88)
            for i, (head, body) in enumerate([("Jose", "Product vision | Brazilian market access | Operator discovery"), ("Meta partner", "Enterprise credibility | Security network leverage | Strategic distribution"), ("Operating model", "Demo first | Benchmark visible | Pilot disciplined")]):
                pdf_card(c, 60 + i * 290, 205, 245, 116, head, body, "role")
            pdf_bullets(c, spec.bullets, 72, 150, 13, 98)
        elif spec.kind == "roadmap":
            pdf_text(c, spec.title, 58, 420, 31, "white", "Helvetica-Bold")
            pdf_text(c, spec.subtitle, 60, 362, 15, "muted", max_chars=94)
            for i, (head, body) in enumerate([("Q2 2026", "Public demo | BR synthetic fintech | SecRL smoke run"), ("Q3 2026", "Full SecRL run | SOC connectors | First BR pilot"), ("Q4 2026", "Compliance hardening | US finance beachhead | Partner packaging"), ("H1 2027", "Global pilots | Channel motion | Enterprise controls")]):
                pdf_card(c, 60 + i * 220, 200, 185, 120, head, body, "milestone gate")
            pdf_text(c, "Rule: every expansion step ships with a repeatable demo, benchmark evidence, and operational runbook.", 72, 92, 12, "gold", "Helvetica-Bold")
        elif spec.kind == "ask":
            pdf_text(c, spec.title, 58, 420, 33, "white", "Helvetica-Bold")
            pdf_text(c, spec.subtitle, 60, 362, 15, "muted", max_chars=92)
            pdf_bullets(c, spec.bullets, 72, 300, 16, 90)
            for i, m in enumerate([("Week 1", "alignment", ""), ("Day 30", "diligence", ""), ("Day 60", "pilot decision", "")]):
                pdf_card(c, 75 + i * 280, 80, 230, 78, *m)
        elif spec.kind == "close":
            pdf_text(c, spec.title, 60, 380, 36, "white", "Helvetica-Bold", max_chars=42)
            pdf_text(c, spec.subtitle, 62, 300, 17, "muted", max_chars=75)
            pdf_bullets(c, spec.bullets, 72, 230, 16, 74)
            pdf_text(c, "quarry.12brain.org", 72, 98, 26, "red2", "Helvetica-Bold")
            c.setFillColor(hx(COLORS["red"]))
            c.roundRect(740, 170, 90, 90, 18, fill=1, stroke=0)
            pdf_text(c, "Q", 770, 205, 42, "white", "Helvetica-Bold")
        c.showPage()
    c.save()


def write_sources():
    source_doc = f"""# Quarry Pitch Deck v1 - Source Notes

Generated: {datetime.now(timezone.utc).isoformat()}

## External sources

1. Anthropic Red Team, "Assessing Claude Mythos Preview's cybersecurity capabilities", April 7, 2026.
   URL: https://red.anthropic.com/2026/mythos-preview/
   Used for: Mythos launch timing, Project Glasswing framing, and why the deck treats Mythos as a cybersecurity market trigger.

2. Federal Reserve Board, Michelle W. Bowman, "Artificial Intelligence in the Financial System", May 1, 2026.
   URL: https://www.federalreserve.gov/newsevents/speech/bowman20260501a.htm
   Used for: Fed framing of Mythos as rapidly evolving AI cybersecurity capability and the Powell/Bessent bank convening reference.

3. U.S. Department of the Treasury, "Treasury Announces Public-Private Initiative to Strengthen Cybersecurity and Risk Management for AI", February 18, 2026.
   URL: https://home.treasury.gov/news/press-releases/sb0395
   Used for: 2026 public-private AI cybersecurity resources and financial-services risk-management framing.

4. U.S. Department of the Treasury, "Managing Artificial Intelligence-Specific Cybersecurity Risks in the Financial Services Sector", press release March 27, 2024.
   URL: https://home.treasury.gov/news/press-releases/jy2212
   Used for: AI-specific operational risk, cybersecurity, fraud, and capability-gap framing.

5. Microsoft SecRL GitHub repository, "Benchmarking LLM agents on Cyber Threat Investigation".
   URL: https://github.com/microsoft/SecRL
   Used for: benchmark identity and threat-investigation benchmark framing.

6. Banco Central do Brasil, Pix statistics pages.
   URLs:
   - https://www.bcb.gov.br/estabilidadefinanceira/pix-em-numeros-estatisticas
   - https://www.bcb.gov.br/estabilidadefinanceira/pix/estatisticas
   Used for: +170M individual Pix users, +7B Pix transactions in January 2026, +R$3T monthly volume in October 2025, and one-day transaction record context.

## Internal Quarry sources

1. {SECRL_MD.relative_to(ROOT)}
   Used for: CARD-012 SecRL pitch smoke-run metrics. Conservative claim only: 8 scenarios, one question per incident, 8/8 correct, average {SECRL['avg_time']:.2f}s, {SECRL['total_tokens']:,} tokens, ${SECRL['total_cost']:.4f} estimated cost.

2. {SECRL_JSON.relative_to(ROOT)}
   Used for: machine-readable SecRL run details.

3. {SCREENSHOT.relative_to(ROOT)}
   Used for: demo screenshot in slide 6.

## Claim controls

- The deck intentionally says "SecRL smoke run" rather than "full official SecRL benchmark", because the local result file documents one question per incident.
- The 6-12 month adoption window is an internal strategic thesis. Regulatory facts supporting urgency are sourced separately from Treasury and Federal Reserve materials.
- No real fintech customer data is disclosed; FinPlay Pagamentos is the fictitious CARD-009 persona.
"""
    (OUT_DIR / "quarry-v1-sources.md").write_text(source_doc)

    content_doc = "\n\n".join(
        f"## Slide {i:02d}: {s.title}\n\n{s.subtitle}\n\n" +
        ("\n".join(f"- {b}" for b in s.bullets) if s.bullets else "") +
        ("\n\nMetrics:\n" + "\n".join(f"- {a}: {b} ({c})" for a, b, c in s.metrics) if s.metrics else "") +
        f"\n\nSource footer: {s.source}"
        for i, s in enumerate(SLIDES, start=1)
    )
    (OUT_DIR / "quarry-v1-content.md").write_text("# Quarry Pitch Deck v1 - Slide Content\n\n" + content_doc)


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx = OUT_DIR / "quarry-v1.pptx"
    pdf = OUT_DIR / "quarry-v1.pdf"
    create_pptx(pptx)
    create_pdf(pdf)
    write_sources()

    timestamp = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")
    backup_pptx = OUT_DIR / f"quarry-v1-backup-{timestamp}.pptx"
    backup_pdf = OUT_DIR / f"quarry-v1-backup-{timestamp}.pdf"
    shutil.copy2(pptx, backup_pptx)
    shutil.copy2(pdf, backup_pdf)

    print(json.dumps({
        "pptx": str(pptx),
        "pdf": str(pdf),
        "backup_pptx": str(backup_pptx),
        "backup_pdf": str(backup_pdf),
        "slides": len(SLIDES),
        "secrl": SECRL,
    }, indent=2))


if __name__ == "__main__":
    main()
